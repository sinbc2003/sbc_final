OUTPUT_PATH = r"E:\0. sbclab\00_sbc_final\scripts\..\data\TeacherFlow_TwoStage.pptx"

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_PATH = "teacherflow_investor_deck.pptx"

# 색상 정의
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
DARK = RGBColor(0x1F, 0x20, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF8, 0xF8, 0xFA)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
SOFT_AMBER = RGBColor(0xFF, 0xF0, 0xD0)
SOFT_GREEN = RGBColor(0xEC, 0xFD, 0xF5)
SOFT_RED = RGBColor(0xFE, 0xF2, 0xF2)
SOFT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE = RGBColor(0xF9, 0x73, 0x16)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_amber_line(slide):
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = AMBER
    shape.line.fill.background()

def add_title_with_divider(slide, title, y=0.5, font_size=32, divider_y=1.3):
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(y), prs.slide_width-Inches(1.6), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = DARK
    tx = slide.shapes.add_shape(
        1, Inches(0.8), Inches(divider_y), Inches(1.5), Inches(0.04)
    )
    tx.fill.solid()
    tx.fill.fore_color.rgb = AMBER
    tx.line.fill.background()

def add_bullet(slide, y, text, color=AMBER, font_size=18):
    # 불릿 사각형
    bullet = slide.shapes.add_shape(
        1, Inches(1.0), Inches(y), Inches(0.12), Inches(0.12)
    )
    bullet.fill.solid()
    bullet.fill.fore_color.rgb = color
    bullet.line.fill.background()
    # 텍스트
    tx = slide.shapes.add_textbox(Inches(1.4), Inches(y-0.03), prs.slide_width-Inches(1.6), Inches(0.45))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = DARK

def add_bullet_list(slide, y_start, lines, amber_idx=None, font_size=18):
    for i, line in enumerate(lines):
        color = AMBER if (amber_idx and i in amber_idx) else DARK
        add_bullet(slide, y_start + i*0.55, line, color=color, font_size=font_size)

def add_centered_text(slide, text, y, font_size, color, bold=False, width=prs.slide_width, height=Inches(0.7)):
    tx = slide.shapes.add_textbox(Inches(0), Inches(y), width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

def add_card(slide, x, y, w, h, bg_color, num, num_color, num_size, label, label_color, label_size):
    card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.fill.background()
    # 숫자
    tx = slide.shapes.add_textbox(Inches(x), Inches(y+0.25), Inches(w), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(num_size)
    p.font.bold = True
    p.font.color.rgb = num_color
    p.alignment = PP_ALIGN.CENTER
    # 라벨
    tx2 = slide.shapes.add_textbox(Inches(x), Inches(y+1.05), Inches(w), Inches(0.35))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = label
    p2.font.size = Pt(label_size)
    p2.font.color.rgb = label_color
    p2.alignment = PP_ALIGN.CENTER

def add_section_slide(bg_color, title, title_size, subtitle, subtitle_size, subtitle_color, y_title, y_sub):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    add_centered_text(slide, title, y_title, title_size, WHITE, bold=True)
    add_centered_text(slide, subtitle, y_sub, subtitle_size, subtitle_color)
    return slide

def add_two_column_feature(slide, title, left_header, left_bg, left_items, right_header, right_bg, right_items):
    add_title_with_divider(slide, title)
    # 좌측 헤더
    lh = slide.shapes.add_shape(1, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.35))
    lh.fill.solid()
    lh.fill.fore_color.rgb = left_bg
    lh.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.35))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_header
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # 우측 헤더
    rh = slide.shapes.add_shape(1, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.35))
    rh.fill.solid()
    rh.fill.fore_color.rgb = right_bg
    rh.line.fill.background()
    tx2 = slide.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.35))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = right_header
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    # 좌측 항목
    for i, item in enumerate(left_items):
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.5+i*0.45), Inches(5.5), Inches(0.35))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY
    # 우측 항목
    for i, item in enumerate(right_items):
        tx = slide.shapes.add_textbox(Inches(7.0), Inches(2.5+i*0.45), Inches(5.5), Inches(0.35))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY

def add_numbered_step(slide, y, num, text, num_color=AMBER, num_size=16, text_size=18):
    # 번호 사각형
    rect = slide.shapes.add_shape(1, Inches(1.0), Inches(y), Inches(0.5), Inches(0.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = num_color
    rect.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(1.0), Inches(y), Inches(0.5), Inches(0.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(num_size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # 텍스트
    tx2 = slide.shapes.add_textbox(Inches(1.6), Inches(y+0.08), prs.slide_width-Inches(2.0), Inches(0.45))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = text
    p2.font.size = Pt(text_size)
    p2.font.color.rgb = DARK

def add_before_after(slide, y, left_label, left_color, left_text, left_text_color, right_label, right_color, right_text, right_text_color):
    # 좌측
    left = slide.shapes.add_shape(1, Inches(1.0), Inches(y), Inches(5.0), Inches(1.2))
    left.fill.solid()
    left.fill.fore_color.rgb = left_color
    left.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(1.0), Inches(y+0.3), Inches(5.0), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_label
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = left_text_color
    p.alignment = PP_ALIGN.CENTER
    # 우측
    right = slide.shapes.add_shape(1, Inches(6.5), Inches(y), Inches(5.8), Inches(1.2))
    right.fill.solid()
    right.fill.fore_color.rgb = right_color
    right.line.fill.background()
    tx2 = slide.shapes.add_textbox(Inches(6.5), Inches(y+0.3), Inches(5.8), Inches(0.7))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = right_label
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = right_text_color
    p2.alignment = PP_ALIGN.CENTER

# ----- 슬라이드 1 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = DARK
# 좌측 세로 앰버 바
bar = slide.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(0.15), Inches(2.0))
bar.fill.solid()
bar.fill.fore_color.rgb = AMBER
bar.line.fill.background()
# TeacherFlow
tx = slide.shapes.add_textbox(Inches(1.3), Inches(1.5), Inches(7), Inches(0.7))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "TeacherFlow"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = WHITE
# 교사 업무 자동화 플랫폼
tx2 = slide.shapes.add_textbox(Inches(1.3), Inches(2.8), Inches(7), Inches(0.7))
tf2 = tx2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "교사 업무 자동화 플랫폼"
p2.font.size = Pt(24)
p2.font.color.rgb = AMBER
# 교실의 시간을 되돌려 드립니다
tx3 = slide.shapes.add_textbox(Inches(1.3), Inches(4.0), Inches(7), Inches(0.45))
tf3 = tx3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = "교실의 시간을 되돌려 드립니다"
p3.font.size = Pt(18)
p3.font.color.rgb = GRAY
# 2026 | Investor Deck
tx4 = slide.shapes.add_textbox(Inches(1.3), Inches(6.0), Inches(7), Inches(0.35))
tf4 = tx4.text_frame
tf4.word_wrap = True
p4 = tf4.paragraphs[0]
p4.text = "2026 | Investor Deck"
p4.font.size = Pt(14)
p4.font.color.rgb = GRAY

# ----- 슬라이드 2 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
# 제목
tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), prs.slide_width-Inches(1.6), Inches(0.7))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "교사는 왜 바쁜가?"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = DARK
# 구분선
divider = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04))
divider.fill.solid()
divider.fill.fore_color.rgb = AMBER
divider.line.fill.background()
# 통계 카드 4개
add_card(slide, 0.8, 2.0, 2.5, 1.8, LGRAY, "52h", AMBER, 36, "주당 평균 근무시간", GRAY, 13)
add_card(slide, 3.8, 2.0, 2.5, 1.8, LGRAY, "40%", RED, 36, "행정 업무 비율", GRAY, 13)
add_card(slide, 6.8, 2.0, 2.5, 1.8, LGRAY, "0개", GRAY, 36, "자동화 도구", GRAY, 13)
add_card(slide, 9.8, 2.0, 2.5, 1.8, LGRAY, "50만", BLUE, 36, "전국 교사 수", GRAY, 13)
# 하단 설명 3줄
desc = [
    "공문 처리, 성적 입력, 양식 작성, 생기부 기록...",
    "교사의 전문성은 수업에 있지만, 시간의 40%를 행정에 뺏기고 있습니다.",
    "대안은? HWP를 다루는 자동화 도구는 사실상 전무합니다."
]
for i, line in enumerate(desc):
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(4.5+i*0.4), prs.slide_width-Inches(1.6), Inches(0.35))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY

# ----- 슬라이드 3 -----
slide = add_section_slide(AMBER, "TeacherFlow", 44, "파일 넣고 → 실행 → 완성", 20, SOFT_AMBER, 2.5, 4.2)

# ----- 슬라이드 4 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
add_title_with_divider(slide, "제품 개요")
bullets = [
    "노코드 업무 자동화 데스크톱 앱",
    "HWP/HWPX, Excel, PDF, DOCX 네이티브 지원",
    "AI 기반 문서 자동 작성 + 양식 채우기",
    "오프라인 동작 (학교 네트워크 제약 대응)",
    "열려있는 한/글·Excel·PPT를 AI가 직접 제어",
    "교사 맞춤 프리셋으로 원클릭 실행"
]
amber_idx = [0,2,4]
for i, line in enumerate(bullets):
    color = AMBER if i in amber_idx else DARK
    add_bullet(slide, 1.8 + i*0.55, line, color=color)

# ----- 슬라이드 5 -----
slide = add_section_slide(AMBER, "핵심 기능", 44, "4가지 핵심 가치", 20, SOFT_AMBER, 2.5, 4.2)

# ----- 슬라이드 6 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
add_two_column_feature(
    slide,
    "① 양식 자동 채우기",
    "기능", DARK,
    ["공문 양식(HWP/Excel) 업로드", "AI가 빈칸 자동 감지", "맥락 기반 내용 생성", "원본 서식 100% 보존", "매크로·수식 영향 없음"],
    "효과", AMBER,
    ["교사 작업: 파일 1개 드래그", "소요 시간: 30분 → 2분", "지원 포맷: HWP, HWPX, XLSX", "COM API + openpyxl 기반", "양식 종류 무제한"]
)

# ----- 슬라이드 7 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
add_two_column_feature(
    slide,
    "② 라이브 문서 제어",
    "기능", DARK,
    ["한/글 Excel PPT 실시간 연결", "열려있는 문서에 AI가 직접 작성", "채팅으로 지시 → 즉시 반영", "25개 COM API 액션"],
    "기술", AMBER,
    ["앱별 LLM 스킬 프롬프트 내장", "GPT-4.1로도 정확한 제어", "프로세스 자동 감지 (10초)", "한/글 COM + 보안 모듈"]
)

# ----- 슬라이드 8 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
add_two_column_feature(
    slide,
    "③ 성적 분석 자동화",
    "분석 항목", DARK,
    ["성적 Excel 업로드", "문항별 정답률 자동 산출", "반별 평균·표준편차", "학생 순위 (상위 N명)"],
    "특징", AMBER,
    ["수식 기반 (값이 아닌 수식 삽입)", "원본 데이터 손상 없음", "새 시트에 분석 결과 생성", "COM API로 실시간 작성"]
)

# ----- 슬라이드 9 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
add_title_with_divider(slide, "④ 노드 기반 워크플로우 (고급)")
bullets = [
    "고급 사용자를 위한 VBA 모드",
    "31개 노드 드래그&드롭 조합",
    "PDF→변환→AI→출력 자유 파이프라인",
    "프리셋으로 저장/공유",
    "엔진이 DAG 위상정렬 → 순차 실행",
    "React Flow 비주얼 에디터"
]
amber_idx = [0,2,4]
for i, line in enumerate(bullets):
    color = AMBER if i in amber_idx else DARK
    add_bullet(slide, 1.8 + i*0.55, line, color=color)

# ----- 슬라이드 10 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
add_title_with_divider(slide, "기술 아키텍처")
bullets = [
    "Tauri v2 — 설치 10MB, 메모리 30MB",
    "Python FastAPI — 33+ API",
    "React + React Flow — 4모드 UI",
    "COM API — HWP/Excel/PPT 제어",
    "LLM 통합 — Claude, GPT, Gemini, llama.cpp",
    "오프라인 우선"
]
amber_idx = [0,2,4]
for i, line in enumerate(bullets):
    color = AMBER if i in amber_idx else DARK
    add_bullet(slide, 1.8 + i*0.55, line, color=color)

# ----- 슬라이드 11 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
add_title_with_divider(slide, "경쟁 환경")
bars = [
    ("ChatGPT/Gemini", RED, "HWP 미지원, 온라인만"),
    ("Inline AI", ORANGE, "HWP만, 양식 채우기 없음"),
    ("MS Copilot", GRAY, "HWP 미지원, 유료"),
    ("TeacherFlow", GREEN, "모든 포맷 + 오프라인 + 무료")
]
for i, (name, color, desc) in enumerate(bars):
    y = 1.8 + i*1.2
    bg = SOFT_GREEN if i==3 else LGRAY
    bar = slide.shapes.add_shape(1, Inches(0.8), Inches(y), Inches(11.5), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bg
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(1.0), Inches(y+0.15), Inches(4.0), Inches(0.45))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    tx2 = slide.shapes.add_textbox(Inches(1.0), Inches(y+0.55), Inches(10.5), Inches(0.35))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(14)
    p2.font.color.rgb = GRAY

# ----- 슬라이드 12 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
add_title_with_divider(slide, "시나리오: 공문 처리")
steps = [
    "교육청에서 공문 도착 (ODT + 첨부)",
    "첨부된 양식(HWP) 드래그&드롭",
    "AI가 빈칸 감지 + 내용 자동 생성",
    "교사 확인 → 즉시 제출"
]
for i, step in enumerate(steps):
    add_numbered_step(slide, 1.8 + i*0.7, i+1, step)
add_before_after(
    slide, 5.2,
    "Before: 30분", SOFT_RED, "Before: 30분", RED,
    "After: 2분", SOFT_GREEN, "After: 2분", GREEN
)

# ----- 슬라이드 13 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
add_title_with_divider(slide, "시나리오: 시험 성적 처리")
steps = [
    "성적 Excel 업로드",
    "'분석 탭 만들어줘' 클릭",
    "문항별 정답률+반별 평균+순위",
    "분석 시트 자동 생성→인쇄"
]
for i, step in enumerate(steps):
    add_numbered_step(slide, 1.8 + i*0.7, i+1, step)
add_before_after(
    slide, 5.2,
    "Before: 2시간", SOFT_RED, "Before: 2시간", RED,
    "After: 30초", SOFT_GREEN, "After: 30초", GREEN
)

# ----- 슬라이드 14 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
add_title_with_divider(slide, "시나리오: 생기부 작성")
steps = [
    "학생 활동 자료 업로드",
    "AI가 생기부 초안 생성",
    "교사가 수정/보완",
    "HWP로 저장→NEIS 붙여넣기"
]
for i, step in enumerate(steps):
    add_numbered_step(slide, 1.8 + i*0.7, i+1, step)
add_before_after(
    slide, 5.2,
    "Before: 40분/명", SOFT_RED, "Before: 40분/명", RED,
    "After: 5분/명", SOFT_GREEN, "After: 5분/명", GREEN
)

# ----- 슬라이드 15 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
add_title_with_divider(slide, "비즈니스 모델")
cards = [
    ("Free", GRAY, "무료", GRAY, "오프라인 기능 전체"),
    ("Pro", AMBER, "₩9,900/월", AMBER, "API AI + 클라우드 동기화"),
    ("School", DARK, "₩100만/년", DARK, "학교 라이센스 + 관리자 대시보드")
]
for i, (tier, head_color, price, price_color, desc) in enumerate(cards):
    x = 0.8 + i*4.1
    card = slide.shapes.add_shape(1, Inches(x), Inches(2.0), Inches(3.7), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = LGRAY
    card.line.fill.background()
    # 헤더
    head = slide.shapes.add_shape(1, Inches(x), Inches(2.0), Inches(3.7), Inches(0.8))
    head.fill.solid()
    head.fill.fore_color.rgb = head_color
    head.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(x), Inches(2.0), Inches(3.7), Inches(0.8))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = tier
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # 가격
    tx2 = slide.shapes.add_textbox(Inches(x), Inches(2.95), Inches(3.7), Inches(0.7))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = price
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = price_color
    p2.alignment = PP_ALIGN.CENTER
    # 설명
    tx3 = slide.shapes.add_textbox(Inches(x), Inches(3.7), Inches(3.7), Inches(0.6))
    tf3 = tx3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = desc
    p3.font.size = Pt(13)
    p3.font.color.rgb = GRAY
    p3.alignment = PP_ALIGN.CENTER

# ----- 슬라이드 16 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
add_title_with_divider(slide, "개발 로드맵")
bullets = [
    "2026 Q2 — 베타 출시 (교사 100명)",
    "2026 Q3 — 정식 출시 + 노드 50개 + 마켓플레이스",
    "2026 Q4 — 학교 라이센스 + 교육청 파일럿",
    "2027 Q1 — 로컬 LLM 완전 오프라인",
    "2027 Q2 — 모바일 동반자 앱 + API 개방"
]
amber_idx = [0,2]
for i, line in enumerate(bullets):
    color = AMBER if i in amber_idx else DARK
    add_bullet(slide, 1.8 + i*0.55, line, color=color)

# ----- 슬라이드 17 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
add_title_with_divider(slide, "팀")
card = slide.shapes.add_shape(1, Inches(2), Inches(2.0), Inches(9), Inches(4.0))
card.fill.solid()
card.fill.fore_color.rgb = LGRAY
card.line.fill.background()
# 신병철
tx = slide.shapes.add_textbox(Inches(2.5), Inches(2.3), Inches(8.0), Inches(0.7))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "신병철"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK
# 창업자/개발자/교사
tx2 = slide.shapes.add_textbox(Inches(2.5), Inches(3.0), Inches(8.0), Inches(0.45))
tf2 = tx2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "창업자 / 개발자 / 현직 고등학교 수학 교사"
p2.font.size = Pt(16)
p2.font.color.rgb = AMBER
# 4줄 설명
desc = [
    '미션: "공교육의 반격을 통한 사교육 시장의 붕괴"',
    "edu-shin.com 운영 (RAG 기반 생기부 시스템)",
    "6대 장비 클러스터 직접 구축·운영",
    "Python/React/COM API 풀스택"
]
for i, line in enumerate(desc):
    tx3 = slide.shapes.add_textbox(Inches(2.5), Inches(3.7+i*0.35), Inches(8.0), Inches(0.35))
    tf3 = tx3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = line
    p3.font.size = Pt(15)
    p3.font.color.rgb = GRAY

# ----- 슬라이드 18 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE
add_amber_line(slide)
add_title_with_divider(slide, "현재 진행 상황")
stats = [
    ("31+", "노드"), ("33+", "API"), ("3개", "라이브제어"),
    ("4모드", "UI"), ("25", "COM액션"), ("3개", "LLM스킬")
]
for i, (num, label) in enumerate(stats):
    x = 0.8 + (i%3)*2.1
    y = 2.0 if i<3 else 4.2
    bg = AMBER if i%2==0 else BLUE
    card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.5), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = LGRAY
    card.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(x), Inches(y+0.3), Inches(2.5), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = bg
    p.alignment = PP_ALIGN.CENTER
    tx2 = slide.shapes.add_textbox(Inches(x), Inches(y+1.05), Inches(2.5), Inches(0.35))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = label
    p2.font.size = Pt(13)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER

# ----- 슬라이드 19 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = DARK
# "투자 제안"
tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), prs.slide_width-Inches(1.6), Inches(0.7))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "투자 제안"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
# 구분선
divider = slide.shapes.add_shape(1, Inches(0.8), Inches(1.6), Inches(1.5), Inches(0.04))
divider.fill.solid()
divider.fill.fore_color.rgb = AMBER
divider.line.fill.background()
# "시드 라운드: 3억원"
tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), prs.slide_width-Inches(1.6), Inches(0.7))
tf2 = tx2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "시드 라운드: 3억원"
p2.font.size = Pt(28)
p2.font.bold = True
p2.font.color.rgb = AMBER
# 4개 항목
items = [
    "풀타임 개발자 2명 (18개월)",
    "서버/인프라 (LLM API 비용)",
    "교사 커뮤니티 구축 + 바이럴 마케팅",
    "교육청 파일럿 운영 (3개 시·도)"
]
for i, item in enumerate(items):
    y = 3.5 + i*0.6
    rect = slide.shapes.add_shape(1, Inches(1.0), Inches(y), Inches(0.4), Inches(0.4))
    rect.fill.solid()
    rect.fill.fore_color.rgb = AMBER
    rect.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(1.0), Inches(y), Inches(0.4), Inches(0.4))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(i+1)
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tx2 = slide.shapes.add_textbox(Inches(1.5), Inches(y+0.03), prs.slide_width-Inches(2.0), Inches(0.45))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = item
    p2.font.size = Pt(18)
    p2.font.color.rgb = SOFT_GRAY

# ----- 슬라이드 20 -----
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = DARK
# 좌측 세로 앰버 바
bar = slide.shapes.add_shape(1, Inches(0.8), Inches(2.0), Inches(0.15), Inches(2.5))
bar.fill.solid()
bar.fill.fore_color.rgb = AMBER
bar.line.fill.background()
# "교실의 시간을\n되돌립니다"
tx = slide.shapes.add_textbox(Inches(1.3), Inches(2.0), Inches(10), Inches(1.4))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "교실의 시간을\n되돌립니다"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
# "TeacherFlow"
tx2 = slide.shapes.add_textbox(Inches(1.3), Inches(4.0), Inches(10), Inches(0.7))
tf2 = tx2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "TeacherFlow"
p2.font.size = Pt(24)
p2.font.bold = True
p2.font.color.rgb = AMBER
# "신병철 | shin@edu-shin.com | edu-shin.com"
tx3 = slide.shapes.add_textbox(Inches(1.3), Inches(5.0), Inches(10), Inches(0.45))
tf3 = tx3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = "신병철 | shin@edu-shin.com | edu-shin.com"
p3.font.size = Pt(16)
p3.font.color.rgb = GRAY
# "감사합니다"
tx4 = slide.shapes.add_textbox(Inches(1.3), Inches(5.4), Inches(10), Inches(0.45))
tf4 = tx4.text_frame
tf4.word_wrap = True
p4 = tf4.paragraphs[0]
p4.text = "감사합니다"
p4.font.size = Pt(16)
p4.font.color.rgb = GRAY

prs.save(OUTPUT_PATH)
print(f"완료: {len(prs.slides)}장")
