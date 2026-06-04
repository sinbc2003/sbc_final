"""
Claude가 직접 만드는 TeacherFlow 마케팅 PPT (python-pptx).
디자인: 앰버/다크 그라디언트, 깔끔한 기업 프레젠테이션.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── 색상 팔레트 ──
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)  # 주 색상
DARK     = RGBColor(0x1F, 0x20, 0x37)  # 배경/제목
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xF8, 0xF8, 0xFA)
GRAY     = RGBColor(0x6B, 0x72, 0x80)
DGRAY    = RGBColor(0x37, 0x41, 0x51)
GREEN    = RGBColor(0x10, 0xB9, 0x81)
BLUE     = RGBColor(0x3B, 0x82, 0xF6)
RED      = RGBColor(0xEF, 0x44, 0x44)

W = Inches(13.333)  # 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ── 헬퍼 ──

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, text="", font_size=12, font_color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = align
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=DGRAY, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, lines, size=16, color=DGRAY, line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(size * (line_spacing - 1))
    return txBox

def section_slide(title, subtitle=""):
    """앰버 섹션 구분 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, AMBER)
    add_text(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5), title, size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(1), subtitle, size=20, color=RGBColor(0xFF, 0xF0, 0xD0), align=PP_ALIGN.CENTER)
    return slide

def content_slide(title, bullets, accent_idx=None):
    """표준 내용 슬라이드: 좌측 제목 + 우측 내용"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    # 상단 라인
    add_rect(slide, 0, 0, W, Inches(0.06), AMBER)
    # 제목
    add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), title, size=32, color=DARK, bold=True)
    # 구분선
    add_rect(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), AMBER)
    # 불릿
    for i, bullet in enumerate(bullets):
        y = Inches(1.8) + Inches(0.55) * i
        # 불릿 포인트
        dot_color = AMBER if accent_idx and i in accent_idx else GRAY
        add_rect(slide, Inches(1.0), y + Inches(0.08), Inches(0.12), Inches(0.12), dot_color)
        add_text(slide, Inches(1.4), y, Inches(10.5), Inches(0.5), bullet, size=18, color=DGRAY)
    return slide

def two_col_slide(title, left_items, right_items, left_title="", right_title=""):
    """2단 레이아웃"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, 0, 0, W, Inches(0.06), AMBER)
    add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), title, size=32, color=DARK, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), AMBER)

    if left_title:
        add_rect(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.5), DARK, left_title, 14, WHITE, True, PP_ALIGN.CENTER)
    if right_title:
        add_rect(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.5), AMBER, right_title, 14, WHITE, True, PP_ALIGN.CENTER)

    y_start = Inches(2.5) if left_title else Inches(1.8)
    for i, item in enumerate(left_items):
        add_text(slide, Inches(1.0), y_start + Inches(0.45) * i, Inches(5.0), Inches(0.4), f"  {item}", size=15, color=DGRAY)
    for i, item in enumerate(right_items):
        add_text(slide, Inches(7.2), y_start + Inches(0.45) * i, Inches(5.0), Inches(0.4), f"  {item}", size=15, color=DGRAY)
    return slide

def stat_card(slide, x, y, number, label, color=AMBER):
    add_rect(slide, x, y, Inches(2.5), Inches(1.8), LGRAY)
    add_text(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.1), Inches(0.8), number, size=36, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.1), Inches(0.5), label, size=13, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
#  슬라이드 작성
# ═══════════════════════════════════════════

# 1. 표지
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
add_rect(s, Inches(0.8), Inches(1.5), Inches(0.15), Inches(2), AMBER)
add_text(s, Inches(1.3), Inches(1.5), Inches(10), Inches(1.2), "TeacherFlow", size=56, color=WHITE, bold=True)
add_text(s, Inches(1.3), Inches(2.8), Inches(10), Inches(0.6), "교사 업무 자동화 플랫폼", size=24, color=AMBER)
add_text(s, Inches(1.3), Inches(4.0), Inches(8), Inches(0.5), "교실의 시간을 되돌려 드립니다", size=18, color=GRAY)
add_text(s, Inches(1.3), Inches(6.0), Inches(8), Inches(0.4), "2026  |  Investor Deck", size=14, color=GRAY)

# 2. 문제 정의
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, W, Inches(0.06), AMBER)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "교사는 왜 바쁜가?", size=36, color=DARK, bold=True)
add_rect(s, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), AMBER)
stat_card(s, Inches(0.8), Inches(2.0), "52h", "주당 평균 근무시간")
stat_card(s, Inches(3.8), Inches(2.0), "40%", "행정 업무 비율", RED)
stat_card(s, Inches(6.8), Inches(2.0), "0개", "자동화 도구", GRAY)
stat_card(s, Inches(9.8), Inches(2.0), "50만", "전국 교사 수", BLUE)
add_multiline(s, Inches(0.8), Inches(4.5), Inches(11), Inches(2), [
    "공문 처리, 성적 입력, 양식 작성, 생기부 기록...",
    "교사의 전문성은 수업에 있지만, 시간의 40%를 행정에 뺏기고 있습니다.",
    "대안은? 거의 없습니다. HWP를 다루는 자동화 도구는 사실상 전무합니다.",
], size=16, color=GRAY)

# 3. 솔루션 개요
section_slide("TeacherFlow", "파일 넣고 → 실행 → 완성")

# 4. 제품 소개
content_slide("제품 개요", [
    "노코드 업무 자동화 데스크톱 앱",
    "HWP/HWPX, Excel, PDF, DOCX 네이티브 지원",
    "AI 기반 문서 자동 작성 + 양식 채우기",
    "오프라인 동작 (학교 네트워크 제약 대응)",
    "열려있는 한/글·Excel·PPT를 AI가 직접 제어",
    "교사 맞춤 프리셋으로 원클릭 실행",
], accent_idx={0, 2, 4})

# 5. 핵심 기능 섹션
section_slide("핵심 기능", "4가지 핵심 가치")

# 6. 기능1: 양식 채우기
two_col_slide("① 양식 자동 채우기",
    ["공문 양식(HWP/Excel) 업로드", "AI가 빈칸 자동 감지", "맥락 기반 내용 생성", "원본 서식 100% 보존", "매크로·수식 영향 없음"],
    ["교사 작업: 파일 1개 드래그", "소요 시간: 30분 → 2분", "지원 포맷: HWP, HWPX, XLSX", "COM API + openpyxl 기반", "양식 종류 무제한"],
    "기능", "효과")

# 7. 기능2: 라이브 제어
two_col_slide("② 라이브 문서 제어",
    ["한/글, Excel, PPT 실시간 연결", "열려있는 문서에 AI가 직접 작성", "채팅으로 지시 → 즉시 반영", "25개 COM API 액션 지원"],
    ["앱별 LLM 스킬 프롬프트 내장", "GPT-4.1로도 정확한 제어", "프로세스 자동 감지 (10초)", "한/글 COM + 보안 모듈"],
    "기능", "기술")

# 8. 기능3: 성적 분석
two_col_slide("③ 성적 분석 자동화",
    ["성적 Excel 업로드", "문항별 정답률 자동 산출", "반별 평균·표준편차", "학생 순위 (상위 N명)"],
    ["수식 기반 (값이 아닌 수식 삽입)", "원본 데이터 손상 없음", "새 시트에 분석 결과 생성", "COM API로 실시간 작성 가능"],
    "분석 항목", "특징")

# 9. 기능4: 노드 에디터
content_slide("④ 노드 기반 워크플로우 (고급)", [
    "고급 사용자를 위한 VBA 모드",
    "31개 노드 드래그&드롭 조합",
    "PDF→변환→AI→출력 자유 파이프라인",
    "프리셋으로 저장/공유 가능",
    "엔진이 DAG 위상정렬 → 순차 실행",
    "React Flow 기반 비주얼 에디터",
])

# 10. 아키텍처
content_slide("기술 아키텍처", [
    "Tauri v2 데스크톱 앱 — 설치 10MB, 메모리 30MB",
    "Python FastAPI 엔진 — 33+ API 엔드포인트",
    "React + React Flow — 4모드 UI (홈/설계/채팅/관리)",
    "COM API — HWP/Excel/PPT 네이티브 제어",
    "LLM 통합 — Claude, GPT, Gemini, 로컬 llama.cpp",
    "오프라인 우선 — 로컬 LLM으로 인터넷 없이 동작",
])

# 11. 경쟁 분석
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, W, Inches(0.06), AMBER)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "경쟁 환경", size=32, color=DARK, bold=True)
add_rect(s, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), AMBER)
competitors = [
    ("ChatGPT / Gemini", "HWP 미지원, 온라인만", RED),
    ("Inline AI", "HWP만, 양식 채우기 없음", RGBColor(0xF9, 0x73, 0x16)),
    ("MS Copilot", "HWP 미지원, 유료", RGBColor(0x84, 0x84, 0x84)),
    ("TeacherFlow", "모든 포맷 + 오프라인 + 무료", GREEN),
]
for i, (name, desc, color) in enumerate(competitors):
    y = Inches(1.8) + Inches(1.2) * i
    add_rect(s, Inches(0.8), y, Inches(11.5), Inches(1.0), LGRAY if i < 3 else RGBColor(0xEC, 0xFD, 0xF5))
    add_text(s, Inches(1.2), y + Inches(0.1), Inches(3.5), Inches(0.4), name, size=18, color=color, bold=True)
    add_text(s, Inches(1.2), y + Inches(0.5), Inches(10), Inches(0.4), desc, size=14, color=GRAY)

# 12~14. 시나리오
for title, steps, time_before, time_after in [
    ("시나리오: 공문 처리", ["교육청에서 공문 도착 (ODT + 첨부)", "첨부된 양식(HWP) 드래그&드롭", "AI가 빈칸 감지 + 내용 자동 생성", "교사 확인 → 즉시 제출"], "30분", "2분"),
    ("시나리오: 시험 성적 처리", ["성적 Excel 파일 업로드", "'분석 탭 만들어줘' 클릭", "문항별 정답률 + 반별 평균 + 순위", "분석 시트 자동 생성 → 인쇄"], "2시간", "30초"),
    ("시나리오: 생기부 작성", ["학생 활동 자료 업로드", "AI가 생기부 초안 자동 생성", "교사가 수정/보완", "HWP로 저장 → NEIS 붙여넣기"], "40분/명", "5분/명"),
]:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    add_rect(s, 0, 0, W, Inches(0.06), AMBER)
    add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), title, size=32, color=DARK, bold=True)
    add_rect(s, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), AMBER)
    for i, step in enumerate(steps):
        y = Inches(1.8) + Inches(0.7) * i
        add_rect(s, Inches(0.8), y, Inches(0.5), Inches(0.5), AMBER, str(i+1), 16, WHITE, True, PP_ALIGN.CENTER)
        add_text(s, Inches(1.6), y + Inches(0.05), Inches(7), Inches(0.5), step, size=18, color=DGRAY)
    # 시간 비교
    add_rect(s, Inches(0.8), Inches(5.2), Inches(5), Inches(1.2), RGBColor(0xFE, 0xF2, 0xF2))
    add_text(s, Inches(1.0), Inches(5.35), Inches(4.5), Inches(0.4), f"Before: {time_before}", size=20, color=RED, bold=True)
    add_rect(s, Inches(6.5), Inches(5.2), Inches(5.8), Inches(1.2), RGBColor(0xEC, 0xFD, 0xF5))
    add_text(s, Inches(6.7), Inches(5.35), Inches(5.3), Inches(0.4), f"After: {time_after}", size=20, color=GREEN, bold=True)

# 15. 비즈니스 모델
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, W, Inches(0.06), AMBER)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "비즈니스 모델", size=32, color=DARK, bold=True)
add_rect(s, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), AMBER)
tiers = [
    ("Free", "무료", "오프라인 기능 전체", GRAY),
    ("Pro", "₩9,900/월", "API 기반 AI + 클라우드 동기화", AMBER),
    ("School", "₩100만/년", "학교 라이센스 + 관리자 대시보드", DARK),
]
for i, (tier, price, desc, color) in enumerate(tiers):
    x = Inches(0.8) + Inches(4.1) * i
    add_rect(s, x, Inches(2.0), Inches(3.7), Inches(3.5), LGRAY)
    add_rect(s, x, Inches(2.0), Inches(3.7), Inches(0.8), color, tier, 20, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), Inches(3.2), Inches(3.1), Inches(0.6), price, size=28, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), Inches(4.0), Inches(3.1), Inches(1.0), desc, size=13, color=GRAY, align=PP_ALIGN.CENTER)

# 16. 로드맵
content_slide("개발 로드맵", [
    "2026 Q2  —  베타 출시 (교사 100명 테스트)",
    "2026 Q3  —  정식 출시 + 노드 50개 + 마켓플레이스",
    "2026 Q4  —  학교 라이센스 + 교육청 파일럿",
    "2027 Q1  —  로컬 LLM 완전 오프라인 전환",
    "2027 Q2  —  모바일 동반자 앱 + API 개방",
], accent_idx={0, 2})

# 17. 팀
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, W, Inches(0.06), AMBER)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "팀", size=32, color=DARK, bold=True)
add_rect(s, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), AMBER)
add_rect(s, Inches(2), Inches(2.0), Inches(9), Inches(4), LGRAY)
add_text(s, Inches(2.5), Inches(2.3), Inches(8), Inches(0.6), "신병철", size=32, color=DARK, bold=True)
add_text(s, Inches(2.5), Inches(3.0), Inches(8), Inches(0.4), "창업자 / 개발자 / 현직 고등학교 수학 교사", size=16, color=AMBER)
add_multiline(s, Inches(2.5), Inches(3.7), Inches(8), Inches(2), [
    '미션: "공교육의 반격을 통한 사교육 시장의 붕괴"',
    "edu-shin.com 운영 (RAG 기반 생기부 시스템)",
    "6대 장비 클러스터 직접 구축·운영",
    "Python/React/COM API 풀스택",
], size=15, color=GRAY)

# 18. 트랙션
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, W, Inches(0.06), AMBER)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), "현재 진행 상황", size=32, color=DARK, bold=True)
add_rect(s, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.04), AMBER)
stats = [("31+", "구현된 노드"), ("33+", "API 엔드포인트"), ("3개", "앱 라이브 제어"), ("4모드", "UI 완성"), ("25", "COM 액션"), ("3개", "LLM 스킬")]
for i, (num, label) in enumerate(stats):
    x = Inches(0.8) + Inches(2.1) * (i % 6)
    y = Inches(2.0) + Inches(2.2) * (i // 6)
    stat_card(s, x, y, num, label, AMBER if i % 2 == 0 else BLUE)

# 19. 투자 제안
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
add_text(s, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8), "투자 제안", size=36, color=WHITE, bold=True)
add_rect(s, Inches(0.8), Inches(1.6), Inches(1.5), Inches(0.04), AMBER)
add_text(s, Inches(0.8), Inches(2.2), Inches(11), Inches(0.8), "시드 라운드: 3억원", size=28, color=AMBER, bold=True)
uses = ["풀타임 개발자 2명 (18개월)", "서버/인프라 (LLM API 비용)", "교사 커뮤니티 구축 + 바이럴 마케팅", "교육청 파일럿 운영 (3개 시·도)"]
for i, use in enumerate(uses):
    y = Inches(3.5) + Inches(0.6) * i
    add_rect(s, Inches(0.8), y, Inches(0.4), Inches(0.4), AMBER, str(i+1), 12, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, Inches(1.5), y, Inches(10), Inches(0.4), use, size=18, color=RGBColor(0xCC, 0xCC, 0xCC))

# 20. 클로징
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
add_rect(s, Inches(0.8), Inches(2.0), Inches(0.15), Inches(2.5), AMBER)
add_text(s, Inches(1.3), Inches(2.0), Inches(10), Inches(1.2), "교실의 시간을\n되돌립니다", size=48, color=WHITE, bold=True)
add_text(s, Inches(1.3), Inches(4.0), Inches(10), Inches(0.5), "TeacherFlow", size=24, color=AMBER, bold=True)
add_multiline(s, Inches(1.3), Inches(5.0), Inches(10), Inches(1.5), [
    "신병철  |  shin@edu-shin.com  |  edu-shin.com",
    "감사합니다",
], size=16, color=GRAY)


# ── 저장 ──
out = os.path.join(os.path.dirname(__file__), "..", "data", "TeacherFlow_Claude.pptx")
prs.save(out)
print(f"저장 완료: {out}")
print(f"슬라이드: {len(prs.slides)}장")
