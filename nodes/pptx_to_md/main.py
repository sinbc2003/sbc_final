"""
PPTX → 마크다운 변환 노드.

python-pptx로 슬라이드별 텍스트를 추출하여
## 헤딩으로 구분된 마크다운을 생성한다.
"""

from pathlib import Path


def execute(inputs: dict, params: dict, context: dict) -> dict:
    pptx_path = inputs["파일"]

    if not Path(pptx_path).exists():
        raise FileNotFoundError(f"PPTX 파일 없음: {pptx_path}")

    from pptx import Presentation

    context["progress"](0.1)
    context["log"]("PPTX 변환 시작")

    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)
    parts = []

    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        texts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        texts.append(line)

            # 표 추출
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    # 헤더 구분선 삽입
                    header = rows[0]
                    col_count = len(table.columns)
                    separator = "| " + " | ".join(["---"] * col_count) + " |"
                    table_md = header + "\n" + separator
                    if len(rows) > 1:
                        table_md += "\n" + "\n".join(rows[1:])
                    texts.append(table_md)

        if texts:
            content = "\n\n".join(texts)
            parts.append(f"## 슬라이드 {slide_num}\n\n{content}")

        if total_slides > 0:
            context["progress"](0.1 + 0.9 * (slide_num / total_slides))

    result = "\n\n---\n\n".join(parts)

    context["progress"](1.0)
    context["log"](f"변환 완료 ({total_slides}장, {len(result)} 글자)")

    return {"텍스트": result}
